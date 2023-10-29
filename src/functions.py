"""
File: src/functions.py
Author: Anthony Schrapffer
Description: File with basic function for img / text replacement in powerpoint (pptx elements).
"""
import pandas as pd
import pptx


def replace_img_slide(slide:pptx.slide, img:pptx.shapes, img_path:str):
    """
    # Replace the picture in the shape object (img) with the image in img_path.

    Args:
        slide (pptx.Presentation.slide): current slide.
        img (pptx.Presentation.slide.shape): current shape in the slide.
        img_path (str): path of the image.
    """
    # Replace the picture in the shape object (img) with the image in img_path.
    imgPic = img._pic
    imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
    imgPart = slide.part.related_part(imgRID)

    with open(img_path, 'rb') as f:
        rImgBlob = f.read()
    # replace
    imgPart._blob = rImgBlob


def replace_text(shp:pptx.shapes,text_in:str,text_out:str):
    """
    Replace the text in the shape object (just replacing a part of the text).

    Args:
        shp (pptx.Presentation.slide.shape): current shape in the slide.
        text_in (str): text to replace.
        text_out (str): Replacing text.
    """    
    p = shp.text_frame.paragraphs[0]
    p_changed = p.text.replace(text_in, text_out)
    change_text(shp,p_changed)


def get_replace_text(shp:pptx.shapes,df_temp:pd.DataFrame) -> str:
    """
    Get the replacing text when multiple change

    Args:
        shp (pptx.Presentation.slide.shape): current shape in the slide.
        df_temp (pd.DataFrame): DataFrame with replacement for this shape only.
        text_out (str): Replacing text.
    """        
    p_changed = shp.text_frame.paragraphs[0].text
    #
    for _,row in df_temp.iterrows():
        text_in, text_out = row["value_in"], row["value_fill"]
        text_in = str(text_in); text_out = str(text_out)
        p_changed = p_changed.replace(text_in, text_out)    
    return p_changed


def change_text(shp:pptx.shapes,text_replace:str, truns:int = 0, runs:bool = False):
    """
    Replace all the text

    Args:
        shp (pptx.Presentation.slide.shape): current shape in the slide.
        text_replace (str): text for replacing the current shape text.
        truns (int): to use if there are different paragraphs.
        runs(bool): to avoid potential repetition of text 
            due to different style (will require further improvement).
    """            
    p = shp.text_frame.paragraphs[truns]
    p.runs[0].text = text_replace
    if runs == 1:
        try:
            if len(shp.text_frame.paragraphs[0].runs)>1:
                for i in range(1, len(shp.text_frame.paragraphs[0].runs)):
                    shp.text_frame.paragraphs[0].runs[i].text = ""
        except:
            pass