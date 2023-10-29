"""
File: src/presentation.py
Author: Anthony Schrapffer
Description: Contains the PresentationCustomized class allowing to
    automatized filling of powerpoint file with given inputs.
"""

from pptx import Presentation
import pandas as pd
import sys
from .functions import replace_img_slide,change_text,replace_text,get_replace_text
import os


infile = "template/template_empty.pptx"


class RaiseError(Exception):
    pass

class PresentationCustomized:
    """
    The PresentationCustomized class will make the change in a powerpoint file based on:

    Attributes:
        infile (str): location of the template file.
        in_data (str): location of the input data.
        dire_portfolio (str): location of the portfolio folder
        out_file (str): location of the output file

    Methods:
        load_inputs(): Load the sheets from the input file.
        test_inputs(): Test if inputs are fine.
        process_slide(page_num): process slide number "page_num" (start by 1).
        replace_image_in_shp(shp): replace image in shape element "shp".
        replace_table_in_shp(shp): replace values shape element "shp" (being a Table).
        replace_text_in_shp(shp): replace text in element "shp".
        process_pptx(): process all the slides.
        save_processed_pptx(): save the output.
    """
    def __init__(self, in_data:str, dire_portfolio:str, out_file:str):
        """
        Initializes and launch the process of powerpoint.

        Args:
            infile (str): location of the template file.
            in_data (str): location of the input data.
            dire_portfolio (str): location of the portfolio folder
            out_file (str): location of the output file
        """
        self.prs = Presentation(infile)
        self.nslides = len(self.prs.slides)
        #
        self.dire_portfolio = dire_portfolio
        self.out_file = out_file
        self.in_data = in_data
        #
        self.load_inputs()
        self.test_inputs()
        #
        self.process_pptx()
        self.save_processed_pptx()

    def load_inputs(self):
        """
        Load the input information.
        """
        dtypes_dict = {'page':int,"element":str,"value_in":str,"value_fill":str}

        def values_conv(x):
            if x == "":
                x = -999
            return int(x)
        

        self.df_input_img = pd.read_excel(self.in_data, sheet_name = "Images")
        self.df_input_data = pd.read_excel(self.in_data, sheet_name = "Data", dtype = dtypes_dict, 
                                           converters={f'table_{c}':values_conv for c in ["i","j"]})

    def test_inputs(self):
        """
        Test the input information.
        """
        # Verification 1
        test = self.df_input_data[self.df_input_data["element"].duplicated()]
        if "None" in test["value_in"]:
            error_message = "Cannot have multiple change for an element, with a full replacement (None in value_in)"
            raise RaiseError(error_message)
        
        # Cannot have repeated j , i combination for a same element
        test = self.df_input_data[self.df_input_data["table_i"] != -999]
        test = test[test[["element","table_i", "table_j"]].duplicated()]
        if test.shape[0]>0:
            error_message = "Cannot have repeated j,i combination for a single cell."
            raise RaiseError(error_message)()    
        
        #Check if image file exists
        for i, row in self.df_input_img.iterrows():
            if not os.path.isfile(self.dire_portfolio + row["dire_figure"]):
                error_message = f"ERROR: image file {self.dire_portfolio + row['dire_figure']} do not exists."
                raise RaiseError(error_message)()      

    def process_slide(self, page_num):
        """
        Process the current slide.

        Args:
            page_num (int): current page number, start by 1.
        """
        self.slide = self.prs.slides[page_num-1]
        shapes = self.slide.shapes

        self.df_input_img_page = self.df_input_img[self.df_input_img["page"] == page_num]
        self.df_input_data_page = self.df_input_data[self.df_input_data["page"] == page_num]

        for shp in self.slide.shapes:
            # IMAGES
            if shp.name in self.df_input_img_page.element.values:
                self.replace_image_in_shp(shp)

            # TEXT 
            elif shp.name in self.df_input_data_page["element"].values:
                # TABLES
                if shp.shape_type == 19: 
                    self.replace_table_in_shp(shp)
                # OTHERS
                else:
                    self.replace_text_in_shp(shp)
                    
    #############################

    def replace_image_in_shp(self,shp):
        """
        Replace image in shape element "shp".

        Args:
            shp (slides.shapes element): shape considered.
        """
        replace_img_slide(self.slide,shp,self.dire_portfolio+self.df_input_img_page[self.df_input_img_page.element == shp.name]["dire_figure"].values[0])

    def replace_table_in_shp(self,shp):
        """
        replace values shape element "shp" (being a Table).

        Args:
            shp (slides.shapes element): shape considered.
        """
        df_temp = self.df_input_data_page[self.df_input_data_page["element"] == shp.name]
        for _,row in df_temp.iterrows():
            i,j = row["table_i"],row["table_j"]
            text_in, text_out = row["value_in"], row["value_fill"]
            #
            if text_in == "None":
                change_text(shp.table.cell(i,j),text_out)
            else:
                replace_text(shp.table.cell(i,j),text_in, text_out)

    def replace_text_in_shp(self, shp):
        """
        Replace text in element "shp".

        Args:
            shp (slides.shapes element): shape considered.
        """
        df_temp = self.df_input_data_page[self.df_input_data_page["element"] == shp.name]
        p_changed = get_replace_text(shp,df_temp)
        change_text(shp,p_changed, runs = True)                

    #############################

    def process_pptx(self):
        """
        Process all the slides.
        """
        for page_num in range(1,self.nslides+1):
            self.process_slide(page_num)

    def save_processed_pptx(self):
        """
        Save the output.
        """
        self.prs.save(self.out_file)

